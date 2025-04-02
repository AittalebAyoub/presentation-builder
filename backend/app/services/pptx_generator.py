from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def generate_powerpoint(sujet, contenu, nom_formateur="AIT TALEB AYOUB", logo_path="ODC_logo.jpeg", output_filename=None):
    """
    Generate a PowerPoint presentation from the content
    
    Args:
        sujet: Title of the presentation
        contenu: Content to include in the presentation
        nom_formateur: Name of the presenter
        logo_path: Path to the logo image
        output_filename: Optional path to save the file (defaults to sujet.pptx)
        
    Returns:
        Path to the generated PPTX file
    """
    # Set default output filename if not provided
    if not output_filename:
        output_filename = f"{sujet}.pptx"
    
    # Initialize presentation
    prs = Presentation()
    slide_height = prs.slide_height
    margin_top = Inches(0.4)
    line_spacing = Inches(0.2)
    left_margin = Inches(0.8)
    content_width = Inches(8)
    
    # Global variables
    current_top = margin_top
    slide = None
    slide_layout = prs.slide_layouts[6]  # Blank layout
    
    # Check if logo exists
    if not os.path.exists(logo_path):
        # Use a default logo or placeholder
        print(f"Warning: Logo not found at {logo_path}")
        logo_path = None
    
    # Functions for space management and content addition
    def check_space_and_add(content_height):
        """Check if there's enough space on the slide and add a new one if needed"""
        nonlocal current_top, slide
        if current_top + content_height > slide_height - Inches(0.5) or slide is None:
            slide = prs.slides.add_slide(slide_layout)
            current_top = margin_top
            
            # Add orange sidebar and logo to each new slide
            orange_color = RGBColor(255, 121, 0)
            left_bar = slide.shapes.add_shape(
                1,  # Rectangle
                0, 0,  # Position x, y
                Inches(0.2), Inches(7.5)  # Width, height
            )
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = orange_color
            left_bar.line.visible = False
            
            # Add logo if it exists
            if logo_path:
                try:
                    slide.shapes.add_picture(
                        logo_path,
                        Inches(0.5), Inches(7),  # Position
                        width=Inches(1.5)  # Fixed width
                    )
                except Exception as e:
                    print(f"Error loading logo: {e}")
                    
        return current_top
    
    def add_title(text, top, size=20, color=RGBColor(255, 121, 0), bold=True):
        """Add a title to the slide"""
        nonlocal slide
        textbox = slide.shapes.add_textbox(left_margin, top, content_width, Inches(0.5))
        p = textbox.text_frame.paragraphs[0]
        p.text = text
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Inter'
        return textbox.height
    
    def add_dynamic_textbox(text, top, max_width=content_width, max_height=Inches(2.5)):
        """Add a text box with dynamic sizing"""
        nonlocal slide
        avg_chars_per_line = 70
        line_height = 0.4  # Line spacing
        num_lines = (len(text) // avg_chars_per_line) + 1
        height = Inches(num_lines * line_height)
        textbox = slide.shapes.add_textbox(left_margin, top, max_width, height)
        
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        p = text_frame.add_paragraph()
        p.text = text
        p.line_spacing = 1.5  # Line spacing
        
        initial_font_size = Pt(13)
        p.runs[0].font.size = initial_font_size
        p.runs[0].font.name = 'Inter'
        
        for _ in range(10):
            if textbox.height <= max_height:
                break
            p.runs[0].font.size -= Pt(1)
            
        return textbox.height
    
    def add_dynamic_bullets(bullets, top, max_width=content_width, max_height=Inches(2.5)):
        """Add bullet points with dynamic sizing"""
        nonlocal slide
        avg_chars_per_line = 50
        line_height = 0.4  # Vertical spacing
        
        total_chars = sum(len(bullet) for bullet in bullets)
        num_lines = (total_chars // avg_chars_per_line) + len(bullets)
        
        height = Inches(num_lines * line_height)
        
        textbox = slide.shapes.add_textbox(left_margin, top, max_width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = f"• {bullet}"
            p.level = 0
            p.line_spacing = 1.5  # Line spacing
            
            initial_font_size = Pt(15)
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = initial_font_size
            run.font.name = 'Inter'
            
        current_font_size = Pt(14)
        for _ in range(10):
            if textbox.height <= max_height:
                break
                
            for paragraph in text_frame.paragraphs:
                if paragraph.runs:
                    paragraph.runs[0].font.size -= Pt(1)
                    
            current_font_size -= Pt(1)
            
        return textbox.height
    
    def add_code_block(code_lines, top):
        """Add a code block to the slide"""
        nonlocal slide
        height = Inches(0.4 + 0.3 * len(code_lines))
        
        left = left_margin
        width = content_width
        
        background_shape = slide.shapes.add_shape(
            1,  # Rectangle
            left,
            top,
            width,
            height
        )
        
        # Background color #FEFEFE
        background_shape.fill.solid()
        background_shape.fill.fore_color.rgb = RGBColor(254, 254, 254)
        
        background_shape.line.visible = True
        background_shape.line.dash_style = 2
        background_shape.line.color.rgb = RGBColor(0, 0, 0)
        background_shape.line.width = Pt(0.9)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        for line in code_lines:
            p = text_frame.add_paragraph()
            p.text = line
            p.level = 0
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(50, 50, 50)
            
        return height
    
    def add_table(data, top):
        """Add a table to the slide"""
        nonlocal slide
        rows, cols = len(data), len(data[0])
        table_width = content_width
        height = Inches(1.5)
        
        table_shape = slide.shapes.add_table(rows, cols, left_margin, top, table_width, height).table
        
        # Header style
        for col_idx in range(cols):
            cell = table_shape.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(133, 179, 222)  # Light blue for header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            
        # Cell style
        for row_idx, row in enumerate(data):
            for col_idx, value in enumerate(row):
                cell = table_shape.cell(row_idx, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                
                # Alternating row colors
                if row_idx > 0:  # Skip header
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(247, 245, 245)  # Light gray
                    
        return height
    
    def add_custom_title_slide(title, subtitle, presenter_name, presenter_title):
        """Add a title slide to the presentation"""
        slide = prs.slides.add_slide(slide_layout)
        
        # Orange sidebar color
        orange_color = RGBColor(255, 121, 0)
        
        # Add orange sidebar
        left_bar = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,  # Position x, y
            Inches(0.2), Inches(7.5)  # Width, height
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = orange_color
        left_bar.line.visible = False
        
        # Add logo if it exists
        if logo_path:
            try:
                slide.shapes.add_picture(
                    logo_path,
                    Inches(8), Inches(0.5),  # Position
                    width=Inches(1.5)  # Fixed width
                )
            except Exception as e:
                print(f"Error loading logo: {e}")
                
        # Main title
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8), Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(60)
        p.font.name = 'Inter'
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 121, 0)
        
        # Subtitle
        textbox = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(8), Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0, 0, 0)
        
    def add_presenter_slide(presenter_name, presenter_title):
        """Add a presenter slide to the presentation"""
        slide = prs.slides.add_slide(slide_layout)
        
        # Orange sidebar color
        orange_color = RGBColor(255, 102, 0)
        
        # Orange sidebar
        left_bar = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,  # Position x, y
            Inches(0.2), Inches(7.5)  # Width, height
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = orange_color
        left_bar.line.visible = False
        
        # Add logo if it exists
        if logo_path:
            try:
                slide.shapes.add_picture(
                    logo_path,
                    Inches(8), Inches(0.5),  # Position
                    width=Inches(1.5)  # Fixed width
                )
            except Exception as e:
                print(f"Error loading logo: {e}")
                
        # Slide title
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8), Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Présenté par"
        p.font.name = 'Inter'
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # Presenter name
# Presenter title
        textbox = slide.shapes.add_textbox(Inches(2.8), Inches(4.5), Inches(8), Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Formateur"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
    
    def add_plan_slide(formation_dates=""):
        """Add a plan slide to the presentation"""
        slide = prs.slides.add_slide(slide_layout)
        
        # Orange sidebar color
        orange_color = RGBColor(255, 102, 0)
        
        # Orange sidebar
        left_bar = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,  # Position x, y
            Inches(0.2), Inches(7.5)  # Width, height
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = orange_color
        left_bar.line.visible = False
        
        # Add logo if it exists
        if logo_path:
            try:
                slide.shapes.add_picture(
                    logo_path,
                    Inches(8), Inches(0.5),  # Position
                    width=Inches(1.5)  # Fixed width
                )
            except Exception as e:
                print(f"Error loading logo: {e}")
                
        # Slide title
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8), Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Plan de formation"
        p.font.size = Pt(60)
        p.font.name = 'Inter'
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # Dates if provided
        if formation_dates:
            textbox = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(8), Inches(1))
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            p.text = formation_dates
            p.font.size = Pt(24)
            p.font.name = 'Inter'
            p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Start generating the presentation
    
    # Add initial slides
    add_custom_title_slide(
        "Formation",
        sujet,
        nom_formateur,
        "Formateur"
    )
    
    add_presenter_slide(
        nom_formateur,
        "Formateur"
    )
    
    add_plan_slide()
    
    # Generate content slides
    for sections in contenu:
        for section in sections:
            # Check space and add a new slide if needed
            title_height = check_space_and_add(Inches(0.6))
            # Add section title
            add_title(section["title"], title_height)
            current_top += Inches(0.6)
            
            for sub in section.get("subsections", []):
                # Check space for the subtitle
                subtitle_height = check_space_and_add(Inches(0.5))
                
                # Add the subtitle
                add_title(f'• {sub["title"]}', subtitle_height, size=16, bold=True, color=RGBColor(0, 0, 0))
                current_top += Inches(0.4) + line_spacing
                
                # Add content by type
                if "content" in sub:
                    content_height = check_space_and_add(Inches(0.6))
                    add_dynamic_textbox(sub["content"], content_height)
                    current_top += Inches(0.6) + line_spacing
                
                if "bullets" in sub:
                    bullets_height = check_space_and_add(Inches(1))
                    dynamic_bullets_height = add_dynamic_bullets(sub["bullets"], bullets_height)
                    current_top += dynamic_bullets_height + line_spacing
                
                if "code" in sub:
                    code_height = check_space_and_add(Inches(0.4 + 0.3 * len(sub["code"])))
                    add_code_block(sub["code"], code_height)
                    current_top += Inches(0.4 + 0.3 * len(sub["code"])) + line_spacing
                
                if "table" in sub:
                    table_height = check_space_and_add(Inches(1.5))
                    add_table(sub["table"], table_height)
                    current_top += Inches(1.5) + line_spacing
    
    # Save the presentation
    prs.save(output_filename)
    return output_filename